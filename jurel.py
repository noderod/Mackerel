"""
BASICS

Provided a topic, finds the information in wikipedia, summarizes it accordingly
to the user's slides number and returns it to the main program.
"""

import wikipedia
from pptx import Presentation
from pptx.util import Inches, Pt
from sumy.parsers.plaintext import PlaintextParser
from sumy.nlp.tokenizers import Tokenizer
from sumy.summarizers.lex_rank import LexRankSummarizer
from sumy.summarizers.lsa import LsaSummarizer
import os, sys
from PIL import Image
import urllib.request
import numpy as np
import random
import gc



# Checks if any of the following elements are present in a list, True if correct
# ltc (list): List to check
# el2c (list): elements to be checked in the list

def anyelem(ltc, el2c):

    for thisel in el2c:
        if thisel in ltc:
            return True
    else:
        return False



################################################################################
# ACTUAL DEVELOPMENT
################################################################################

# Needs at least the topic and the number of pages
if len(sys.argv) < 2:
    raise SyntaxError('Minimum number of arguments is Topic (str)')


postopic = str(sys.argv[1]) # 'Augustus'

postopic = postopic.replace('-', ' ')

# Creates a new for the presentation
if len(sys.argv) == 2:
    pw_title = postopic
    pw_slides = 8 # Assumes the number of slides to be 8
    pw_sub = ''

else:
    pw_ori = ''

    try:
        pw_slides = int(float(sys.argv[2]))
        numisl = 3
    except:
        numisl = 2
        pw_slides = 8
        print('Number of slides not provided, defaults to 8')

    for hh in range(numisl, len(sys.argv)):
        pw_ori += str(sys.argv[hh])+' '

    # Actually creates the real titles and subtitle
    pw_title = pw_ori.split('/+/')[0]

    try:
        pw_sub = pw_ori.split('/+/')[1]
    except:
        pw_sub = ''


wiktop = wikipedia.page(postopic).content


sumtop = wikipedia.summary(postopic, sentences = 1)

# List of words which indicates a certain topic

artictop = {
    'person' : ['activist', 'actor', 'artist', 'author', 'bussinessman', 'CEO', 'designer',
                'doctor', 'emperor', 'entrepeneur', 'inventor', 'king', 'photographer',
                'politician', 'polymath', 'president', 'producer', 'prophet', 'writer'],
    'place' : ['capital', 'city', 'country', 'municipality', 'prefecture', 'state',
              'town', 'village'],
    'art' : ['architect', 'film', 'movie', 'painting', 'statue'],
    'common object' : [],
    'lifeform' : ['animal', 'bacteria', 'organism', 'plant', 'protozoo', 'species'],
    'religion' : ['relig', 'theism'],
    'ideology' : ['politic']
}


# Determines the topic of the article
notfound = True

for postt in artictop.keys():

    for identifier in artictop[postt]:
        if identifier in sumtop.lower():
            topic = postt
            notfound = False
            break

    if notfound == False:
        break
else:
    topic = 'Unknown'


print(topic)


# If the topic is a person, then there are several things needed
if topic == 'person':
    occupations = []
    for posocc in artictop['person']:
        if posocc in sumtop:
            occupations.append(posocc)

    # Obtains the birth-death dates
    repbar = sumtop.replace('(', '|||').replace(')', '|||').split('|||')
    for elem in repbar:
        if ';' in elem:
            birdea = elem
            break

    birdea = birdea.split(';')[1]
    birdea = birdea.split('–') # This is not a keyboard key
    if len(birdea) == 1:
        # Still alive
        Bdat = birdea[0]
        Ddat = 'Alive'
    else:
        Bdat = birdea[0]
        Ddat = birdea[1]


# Dumps all the content into a text file, except the last part
tobeus = wiktop.split('== See also ==')


# Skips the titles
actualtext = tobeus[0].replace('=', '\n')
gc.collect()


# Transforms the unicode into ascii to avoid future problems
tobeus_enc = actualtext.encode('ascii', errors='ignore').decode()

with open('Tempfile___.txt', 'w') as lolfil:
    lolfil.write(tobeus_enc)


thefile = "Tempfile___.txt" #name of the plain-text file
parser = PlaintextParser.from_file(thefile, Tokenizer("english"))
summarizer = LexRankSummarizer()

summary = summarizer(parser.document, 9*(pw_slides-1)-1) #Summarize the document with 5 sentences


# Creates a list of lines that will be used
usedlin = []

for sentence in summary:

    usedlin.append(str(sentence))


# summarizes the text file using the number of slides provided by the user
os.remove('Tempfile___.txt') # Removes the file
gc.collect


# Orders the sentences in a single string according to a biographical order.
# allsen (arr) (str): Conglomerate of all sentences
# Returns a list of the sentences, NOT a single str

def order_sen(allsen):

    fwq = []
    for mem in allsen:
        try:
            datdat = min([int(s) for s in mem.split() if s.isdigit()])

        except:
            # For sentences without numbers in them
            datdat = np.inf

        if 'BC' in mem:
            datdat = -datdat

        fwq.append([mem, datdat])

    wer =  sorted(fwq, key = lambda x: x[1])
    nam = []

    for kol in wer:
        nam.append(kol[0])

    return nam


usedlin = order_sen(usedlin)
################################################################################
# ACTUAL PRESENTATION
################################################################################

# Obtains a series of images
allimg = wikipedia.page(postopic).images


# Reduces a set of parameters in order to obtain a ratio
# s1 (arr): Set to be checked
# baremo (arr): original set
# Returns a ratio (float),  1 if not needed

def ratcalc(s1, baremo):

    rats = [1] # All ratios

    for xx, yy in zip(s1, baremo):

        if xx > yy:
            rats.append(yy/xx)


    return min(rats)


# Locates an image in a slide in the presentation
# goose (str): URL of the image

def im2pres(goose):

    urllib.request.urlretrieve(goose, 'ImgPrin___.jpg')

    # Obviously, the location of the image depends on its size
    with Image.open('ImgPrin___.jpg') as test_image:
        width, heiimg = test_image.size

    # Finds the qualifier that dimishes the size if the image is too big
    # Assumes that the slide is divided as a rectangle of 1080x810 pixels
    ratiopic = ratcalc([width, heiimg], [1080, 810])
    # The important parameters are as follow
    leftdis = 1080 - ratiopic*width
    heidis = ratiopic*heiimg

    pic = slide.shapes.add_picture('ImgPrin___.jpg', Inches(leftdis*(10.4/1080)), Inches(0), height=Inches(heidis*(7.2/810)))
    os.remove('ImgPrin___.jpg')



# Title slide, always constant

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = pw_title
subtitle.text = pw_sub

# In the case of people, the dates of birth and death are also shown
if topic == 'person':
    pw_sub += '\n'+Bdat+' - '+Ddat

subtitle.text = pw_sub

# Image goes first, uses the starting wikipedia image
# Needs to create a temporary image as an intermediary
# Sometimes there are OS errors, if so, the system chooses a different image
img_seen_already = []

while True:

    # No image shown if there are no available images
    if (len(allimg) == 0) or (len(allimg) == len(img_seen_already)):
        break

    imurl = random.choice(allimg)

    if imurl in img_seen_already:
        continue

    img_seen_already.append(imurl)

    try:
        im2pres(imurl)
        # If it works, no need to check any more images
        break
    except:
        pass


# CREATES THE NORMAL SLIDES

actual_slide_n = pw_slides-1 # Because the title does not count
# Each slide has 5 sentences
slide_sentences = []
print(len(usedlin))

for cowcow in range(0, actual_slide_n):
    slide_sentences.append([])

    # Breaks if there are no more sentences
    if (len(usedlin) - (8*len(slide_sentences))) < 8:
        break

    for vtvt in range(0, 8):
        slide_sentences[-1].append(usedlin[8*cowcow + vtvt])


# Summarizes a set of lines into one, preferably for titles
#  sena (str): Contains the sentences

def slititle(sena):

    # Creates a temporary text file since sumy works that way
    with open('Titlefile___.txt', 'w') as tefil:
        tefil.write(sena)


    thefile = "Titlefile___.txt" #name of the plain-text file
    parser = PlaintextParser.from_file(thefile, Tokenizer("english"))
    summarizer = LsaSummarizer()

    summary = summarizer(parser.document, 1) # Reduce the document to 1 sentence
    os.remove('Titlefile___.txt')

    return str(summary[0]).split('.')[0]


# Creates a slide given a set of lines and pictures
# withsent (list) (str): Contains the sentences

def one_more_slide(withsent):

    comptext = ''
    for tyu in range(0, len(withsent)-1):
        comptext += withsent[tyu] + '\n'

    else:
        comptext += withsent[-1]

    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    left = top = width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    TX = txBox.text_frame
    # Title
    TX.text = slititle(comptext)

    # Actual content
    p = TX.add_paragraph()
    p.text = comptext
    p.font.size = Pt(12)


for thissen in slide_sentences:

    if len(thissen) == 0:
        continue

    one_more_slide(order_sen(thissen))


prs.save(postopic+'.pptx')
